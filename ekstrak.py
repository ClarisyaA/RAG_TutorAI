import os
import fitz  # PyMuPDF
from pptx import Presentation
from langdetect import detect
from deep_translator import GoogleTranslator
from docx import Document


def extract_text_from_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_text_from_ppt(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def extract_text_from_docx(path):
    doc = Document(path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def detect_language(text):
    try:
        return detect(text)
    except:
        return "unknown"


def translate_to_english(text):
    return GoogleTranslator(source='id', target='en').translate(text)


def extract_text(filepath):
    ext = os.path.splitext(filepath)[-1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(filepath)
    elif ext == ".pptx":
        return extract_text_from_ppt(filepath)
    elif ext == ".txt":
        return extract_text_from_txt(filepath)
    elif ext == ".docx":
        return extract_text_from_docx(filepath)
    else:
        print(f"❌ Lewatkan file tidak didukung: {filepath}")
        return None


def save_texts(base_name, original, translated=None):
    out_dir = "ekstrak"
    os.makedirs(out_dir, exist_ok=True)

    # Simpan hasil utama
    target_path = os.path.join(out_dir, base_name + ".txt")
    with open(target_path, "w", encoding="utf-8") as f:
        f.write(translated if translated else original)

    # Jika diterjemahkan, simpan versi asli dan versi en
    if translated:
        with open(os.path.join(out_dir, base_name + ".id.txt"), "w", encoding="utf-8") as f:
            f.write(original)
        with open(os.path.join(out_dir, base_name + ".en.txt"), "w", encoding="utf-8") as f:
            f.write(translated)


if __name__ == "__main__":
    materi_dir = "Materi"
    files = os.listdir(materi_dir)

    print(f"📂 Memproses semua file dalam folder '{materi_dir}'...")

    for file in files:
        filepath = os.path.join(materi_dir, file)
        base_name, ext = os.path.splitext(file)
        ext = ext.lower()

        if ext not in [".pdf", ".pptx", ".txt", ".docx"]:
            print(f"⚠️ Lewatkan '{file}' (format tidak didukung)")
            continue

        print(f"🔍 Mengekstrak: {file}")
        raw_text = extract_text(filepath)

        if not raw_text or len(raw_text.strip()) == 0:
            print(f"⚠️ Tidak ada teks di: {file}")
            continue

        lang = detect_language(raw_text)
        print(f"🌐 Deteksi bahasa: {lang}")

        if lang == "id":
            print(f"🔄 Menerjemahkan {file} ke Bahasa Inggris...")
            translated = translate_to_english(raw_text)
            save_texts(base_name, raw_text, translated)
        else:
            save_texts(base_name, raw_text)

    print("✅ Proses ekstraksi selesai. Hasil di folder 'ekstrak/'")
